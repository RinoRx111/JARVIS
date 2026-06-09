import os
import logging
from typing import List, Dict, Any
from sqlmodel import Session, create_engine
from app.models.file import FileMetadata
from app.services.memory import memory_service

logger = logging.getLogger(__name__)

def chunk_text(text: str, chunk_size: int = 1000, overlap: int = 100) -> List[str]:
    """Chunks text into smaller segments for vector embedding."""
    if not text:
        return []
    chunks = []
    start = 0
    text_len = len(text)
    while start < text_len:
        end = min(start + chunk_size, text_len)
        chunks.append(text[start:end])
        if end == text_len:
            break
        start += chunk_size - overlap
    return chunks

def extract_text_from_file(file_path: str) -> str:
    """Parses files of different extensions and returns their raw text contents."""
    filename = os.path.basename(file_path).lower()
    if filename.endswith(".pdf"):
        from pypdf import PdfReader
        reader = PdfReader(file_path)
        text = ""
        for page in reader.pages:
            text += (page.extract_text() or "") + "\n"
        return text
    elif filename.endswith(".docx"):
        import docx
        doc = docx.Document(file_path)
        text = ""
        for paragraph in doc.paragraphs:
            text += paragraph.text + "\n"
        return text
    elif filename.endswith(".csv"):
        import pandas as pd
        df = pd.read_csv(file_path)
        # Represent CSV rows as textual representation
        return df.to_string(index=False)
    elif filename.endswith((".xlsx", ".xls")):
        import pandas as pd
        try:
            df = pd.read_excel(file_path)
            return df.to_string(index=False)
        except Exception as e:
            return f"[Excel file reading error: {e}]"
    elif filename.endswith(".pptx"):
        try:
            import pptx
            prs = pptx.Presentation(file_path)
            text = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        text.append(shape.text)
            return "\n".join(text)
        except Exception as e:
            return f"[PowerPoint presentation parsing error: {e}]"
    elif filename.endswith((".png", ".jpg", ".jpeg", ".bmp", ".webp", ".gif")):
        from PIL import Image
        from app.core.config import settings
        img = Image.open(file_path)
        info = f"Image File: {filename}\nFormat: {img.format}\nSize: {img.size}\nMode: {img.mode}\n"
        
        # Multimodal Vision Analysis fallback
        if settings.GEMINI_API_KEY:
            try:
                import google.generativeai as genai
                genai.configure(api_key=settings.GEMINI_API_KEY)
                model = genai.GenerativeModel('gemini-1.5-flash')
                response = model.generate_content([
                    "Please provide a highly detailed description of this image for an AI assistant. List all text, objects, and visual context you can see.",
                    img
                ])
                info += f"Visual Description:\n{response.text}"
                return info
            except Exception as e:
                logger.error(f"Gemini image description failed: {e}")
                
        if settings.OPENAI_API_KEY:
            try:
                import base64
                from openai import OpenAI
                client = OpenAI(api_key=settings.OPENAI_API_KEY)
                with open(file_path, "rb") as image_file:
                    base64_image = base64.b64encode(image_file.read()).decode('utf-8')
                
                response = client.chat.completions.create(
                    model="gpt-4o",
                    messages=[
                        {
                            "role": "user",
                            "content": [
                                {"type": "text", "text": "Please provide a highly detailed description of this image for an AI assistant. List all text, objects, and visual context you can see."},
                                {
                                    "type": "image_url",
                                    "image_url": {
                                        "url": f"data:image/jpeg;base64,{base64_image}"
                                    }
                                }
                            ]
                        }
                    ]
                )
                info += f"Visual Description:\n{response.choices[0].message.content}"
                return info
            except Exception as e:
                logger.error(f"OpenAI image description failed: {e}")
                
        return info + "\nVisual Description: [Vision API Key not configured. Visual description omitted.]"
    elif filename.endswith(".txt") or filename.endswith(".md"):
        with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
            return f.read()
    else:
        raise ValueError("Unsupported file format for text extraction")

def ingest_file_in_background(db_url: str, file_record_id: int):
    """
    Background worker function that parses a file, chunks it,
    and indexes it in ChromaDB, updating the file status in the database.
    """
    engine = create_engine(db_url)
    with Session(engine) as db:
        # Fetch file record
        file_record = db.get(FileMetadata, file_record_id)
        if not file_record:
            logger.error(f"FileMetadata record with ID {file_record_id} not found.")
            return

        try:
            logger.info(f"Starting asynchronous ingestion for file: {file_record.filename} (ID: {file_record_id})")
            file_record.status = "processing"
            db.commit()
            db.refresh(file_record)

            if not os.path.exists(file_record.file_path):
                raise FileNotFoundError(f"Physical file not found at {file_record.file_path}")

            # Step 1: Extract text
            raw_text = extract_text_from_file(file_record.file_path)
            
            # Step 2: Chunk text
            chunks = chunk_text(raw_text)
            
            # Step 3: Embed and index
            if chunks:
                for idx, chunk in enumerate(chunks):
                    memory_service.add_document_chunk(
                        file_id=str(file_record.id),
                        chunk_index=idx,
                        content=chunk.strip(),
                        metadata={
                            "filename": file_record.filename,
                            "owner_id": file_record.user_id,
                            "file_id": file_record.id
                        }
                    )
            
            logger.info(f"Indexed {len(chunks)} chunks successfully for file ID {file_record_id}")
            file_record.status = "completed"
            db.commit()

        except Exception as e:
            logger.error(f"Failed to ingest file {file_record_id}: {e}", exc_info=True)
            file_record.status = "failed"
            file_record.error = str(e)
            db.commit()
